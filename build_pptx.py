"""
Zambia Economic Recovery — Follow-Up Assignment
Storytelling with Data: Big Idea + 3-Minute Story Structure
Author: Boldwin Mweemba
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
import copy

# ── Chart data ────────────────────────────────────────────────────────────────
YEARS      = ['2015','2016','2017','2018','2019','2020','2021','2022','2023','2024']
EXT_DEBT   = [12.3, 14.0, 17.5, 22.4, 29.1, 31.5, 29.8, 26.4, 22.7, 18.9]  # $B
CURR_ACCT  = [-3.6, -4.1, -1.9, -1.7, -1.5, 2.9, 11.9, 4.7, 3.1, 2.0]      # % GDP
GDP_GROWTH = [2.9, 3.8, 3.5, 4.0, 1.4, -2.8, 4.6, 4.7, 4.7, 5.3]           # % annual
FDI        = [3.1, 2.4, 4.2, 1.7, 2.5, 1.9, 3.6, 4.9, 6.8, 9.32]           # % GDP (World Bank WDI)


WHITE = RGBColor(0xFF, 0xFF, 0xFF)

def add_white_bg(slide, left, top, width, height):
    """White rectangle placed BEFORE a chart so it acts as the chart background."""
    add_rect(slide, left, top, width, height, WHITE)


def add_column_chart(slide, left, top, width, height, categories, values,
                     series_color=None, point_colors=None):
    cd = ChartData()
    cd.categories = categories
    cd.add_series('', values)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, left, top, width, height, cd).chart
    chart.has_legend = False
    series = chart.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = series_color or ACCENT
    series.format.line.color.rgb = series_color or ACCENT
    if point_colors:
        for idx, col in point_colors.items():
            pt = series.points[idx]
            pt.format.fill.solid()
            pt.format.fill.fore_color.rgb = col
    try:
        chart.plots[0].gap_width = 80
        chart.category_axis.tick_labels.font.size = Pt(7)
        chart.category_axis.tick_labels.font.name = "Calibri"
        chart.value_axis.tick_labels.font.size = Pt(7)
        chart.value_axis.tick_labels.font.name = "Calibri"
    except Exception:
        pass
    return chart


def add_line_chart(slide, left, top, width, height, categories, values,
                   series_color=None):
    cd = ChartData()
    cd.categories = categories
    cd.add_series('', values)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE, left, top, width, height, cd).chart
    chart.has_legend = False
    series = chart.series[0]
    series.smooth = True
    series.format.line.width = Pt(2.0)
    series.format.line.color.rgb = series_color or ACCENT2
    try:
        chart.category_axis.tick_labels.font.size = Pt(7)
        chart.category_axis.tick_labels.font.name = "Calibri"
        chart.value_axis.tick_labels.font.size = Pt(7)
        chart.value_axis.tick_labels.font.name = "Calibri"
    except Exception:
        pass
    return chart

# ── Colour palette ────────────────────────────────────────────────────────────
DARK_BG    = RGBColor(0x0D, 0x1B, 0x2A)   # deep navy
ACCENT     = RGBColor(0xE8, 0x8C, 0x00)   # amber / copper tone
ACCENT2    = RGBColor(0x2E, 0xCC, 0x71)   # emerald green (growth)
LIGHT_TEXT = RGBColor(0xF5, 0xF5, 0xF5)   # near-white
MID_TEXT   = RGBColor(0xB0, 0xBE, 0xC5)   # muted grey-blue
RED_WARN   = RGBColor(0xE7, 0x4C, 0x3C)   # warning red
DIVIDER    = RGBColor(0x1E, 0x3A, 0x5F)   # slightly lighter navy

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ── Helper utilities ──────────────────────────────────────────────────────────

def set_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color, alpha=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, font_size, bold=False,
                color=None, align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color if color else LIGHT_TEXT
    run.font.name = "Calibri"
    return txBox


def add_multiline(slide, left, top, width, height, lines, font_size,
                  color=None, bold=False, align=PP_ALIGN.LEFT, line_spacing=None):
    """lines = list of (text, bold_override, color_override) or plain strings."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in lines:
        if isinstance(item, str):
            text, b, c = item, bold, color
        else:
            text, b, c = item[0], item[1] if len(item) > 1 else bold, item[2] if len(item) > 2 else color
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = b
        run.font.color.rgb = c if c else (color if color else LIGHT_TEXT)
        run.font.name = "Calibri"
    return txBox


def add_divider_line(slide, left, top, width, color=None):
    line = slide.shapes.add_shape(1, left, top, width, Pt(1.5))
    line.fill.solid()
    line.fill.fore_color.rgb = color if color else ACCENT
    line.line.fill.background()
    return line


def stat_card(slide, left, top, width, height, metric, value, sub, val_color=None):
    add_rect(slide, left, top, width, height, DIVIDER)
    add_textbox(slide, left + Inches(0.1), top + Inches(0.1),
                width - Inches(0.2), Inches(0.45),
                metric, 14, color=MID_TEXT, align=PP_ALIGN.CENTER)
    add_textbox(slide, left + Inches(0.05), top + Inches(0.52),
                width - Inches(0.1), Inches(0.58),
                value, 24, bold=True,
                color=val_color if val_color else ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, left + Inches(0.1), top + Inches(1.08),
                width - Inches(0.2), Inches(0.42),
                sub, 12, color=MID_TEXT, align=PP_ALIGN.CENTER)


# ── Presentation setup ────────────────────────────────────────────────────────

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]   # completely blank


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title / Big Idea
# ══════════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(blank_layout)
set_bg(s1, DARK_BG)

# copper-tone accent strip on left
add_rect(s1, Inches(0), Inches(0), Inches(0.12), SLIDE_H, ACCENT)

# decorative top band
add_rect(s1, Inches(0.12), Inches(0), SLIDE_W - Inches(0.12), Inches(0.06), DIVIDER)

# country + project label
add_textbox(s1, Inches(0.5), Inches(0.35), Inches(8), Inches(0.4),
            "ZAMBIA  ·  MACROECONOMIC ANALYSIS  ·  2015–2025",
            9, color=ACCENT, bold=True)

# main title
add_multiline(s1, Inches(0.5), Inches(1.0), Inches(12.4), Inches(2.2),
              [("Zambia's Debt Restructuring", True, LIGHT_TEXT),
               ("Converted a Commodity Windfall", True, LIGHT_TEXT),
               ("Into a Durable Recovery", True, ACCENT)],
              font_size=38, align=PP_ALIGN.LEFT)

# divider
add_divider_line(s1, Inches(0.5), Inches(3.2), Inches(9))

# BIG IDEA label
add_textbox(s1, Inches(0.5), Inches(3.35), Inches(1.8), Inches(0.35),
            "BIG IDEA", 9, bold=True, color=ACCENT)

# Big Idea sentence
big_idea = (
    "Zambia's G20 Common Framework debt restructuring converted a copper-driven "
    "fiscal windfall into a durable institutional framework — slashing PPG debt "
    "service from 12.5 % to 3 % of exports, sustaining 5 %+ GDP growth, and "
    "attracting record FDI of 9.32 % of GDP in 2024 — proving that commodity "
    "luck alone cannot anchor a recovery without credible institutional reform."
)
add_textbox(s1, Inches(0.5), Inches(3.7), Inches(11.8), Inches(1.5),
            big_idea, 13, color=LIGHT_TEXT, italic=True)

# bottom author line
add_textbox(s1, Inches(0.5), Inches(6.9), Inches(6), Inches(0.4),
            "Boldwin Mweemba  ·  Data Analyst  ·  2026",
            9, color=MID_TEXT)

# slide number
add_textbox(s1, Inches(12.5), Inches(6.9), Inches(0.7), Inches(0.4),
            "01", 9, color=MID_TEXT, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — The Problem: Zambia's Debt Spiral (2015–2020)
# ══════════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(blank_layout)
set_bg(s2, DARK_BG)
add_rect(s2, Inches(0), Inches(0), Inches(0.12), SLIDE_H, RED_WARN)
add_rect(s2, Inches(0.12), Inches(0), SLIDE_W - Inches(0.12), Inches(0.06), DIVIDER)

add_textbox(s2, Inches(0.5), Inches(0.25), Inches(8), Inches(0.4),
            "3-MINUTE STORY  ·  PART 1 OF 4  ·  THE PROBLEM", 9, color=RED_WARN, bold=True)

add_multiline(s2, Inches(0.5), Inches(0.7), Inches(12), Inches(1.1),
              [("The Debt Spiral: How Zambia Reached Default", True, LIGHT_TEXT)],
              font_size=28)

add_divider_line(s2, Inches(0.5), Inches(1.75), Inches(5.5), color=RED_WARN)

# narrative text
narrative = [
    ("Between 2015 and 2019, Zambia's external debt tripled — from $12.3 B to $29.1 B — "
     "as the government borrowed heavily against future copper revenues. When COVID-19 struck "
     "in 2020, the fiscal position collapsed. In November 2020, Zambia became the first "
     "African sovereign to default during the pandemic era, missing a $42.5 M Eurobond "
     "coupon payment.", False, LIGHT_TEXT),
    ("", False, LIGHT_TEXT),
    ("Two structural vulnerabilities converged:", False, MID_TEXT),
    ("  •  A copper-dependent export base unable to absorb a commodity price shock", False, MID_TEXT),
    ("  •  Debt service obligations that consumed 12.5 % of all export earnings", False, MID_TEXT),
]
add_multiline(s2, Inches(0.5), Inches(2.0), Inches(6.0), Inches(3.0),
              narrative, font_size=12)

# External Debt bar chart (replaces text timeline boxes)
add_textbox(s2, Inches(7.2), Inches(1.65), Inches(5.7), Inches(0.3),
            "EXTERNAL DEBT  ($ billions)  ·  2015–2024", 8, bold=True, color=MID_TEXT)
add_white_bg(s2, Inches(7.2), Inches(1.95), Inches(5.85), Inches(3.1))
add_column_chart(s2, Inches(7.2), Inches(1.95), Inches(5.85), Inches(3.1),
                 YEARS, EXT_DEBT, series_color=ACCENT,
                 point_colors={4: RED_WARN, 5: RED_WARN})

# stat
add_rect(s2, Inches(0.5), Inches(5.3), Inches(5.8), Inches(1.0), DIVIDER)
add_textbox(s2, Inches(0.65), Inches(5.4), Inches(5.5), Inches(0.4),
            "PPG DEBT SERVICE AT PEAK", 9, color=MID_TEXT)
add_textbox(s2, Inches(0.65), Inches(5.75), Inches(2.5), Inches(0.45),
            "12.5%", 22, bold=True, color=RED_WARN)
add_textbox(s2, Inches(3.0), Inches(5.85), Inches(3.0), Inches(0.35),
            "of total export earnings", 11, color=MID_TEXT)

add_textbox(s2, Inches(12.5), Inches(6.9), Inches(0.7), Inches(0.4),
            "02", 9, color=MID_TEXT, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — The Inflection: Copper Supercycle + IMF Programme
# ══════════════════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(blank_layout)
set_bg(s3, DARK_BG)
add_rect(s3, Inches(0), Inches(0), Inches(0.12), SLIDE_H, ACCENT)
add_rect(s3, Inches(0.12), Inches(0), SLIDE_W - Inches(0.12), Inches(0.06), DIVIDER)

add_textbox(s3, Inches(0.5), Inches(0.25), Inches(9), Inches(0.4),
            "3-MINUTE STORY  ·  PART 2 OF 4  ·  THE INFLECTION", 9, color=ACCENT, bold=True)

add_multiline(s3, Inches(0.5), Inches(0.7), Inches(12), Inches(1.1),
              [("Copper Created the Breathing Room", True, LIGHT_TEXT)],
              font_size=28)

add_divider_line(s3, Inches(0.5), Inches(1.75), Inches(5.5))

narrative3 = [
    ("Before the restructuring deal was signed, a copper price supercycle transformed "
     "Zambia's external position. The current account swung from a -3.6 % deficit in "
     "2015 to a +11.9 % surplus in 2021 — purely on the back of commodity revenue.", False, LIGHT_TEXT),
    ("", False, LIGHT_TEXT),
    ("Simultaneously, the IMF approved a $1.3 B Extended Credit Facility (ECF) "
     "programme, giving the government a credible fiscal anchor and unlocking "
     "multilateral support.", False, LIGHT_TEXT),
    ("", False, LIGHT_TEXT),
    ("KEY INSIGHT: The copper windfall was necessary but not sufficient. "
     "Without institutional reform, history suggested it would be spent — not saved.",
     True, ACCENT),
]
add_multiline(s3, Inches(0.5), Inches(2.0), Inches(6.2), Inches(3.5),
              narrative3, font_size=12)

# stat cards
cards = [
    ("Current Account\n2015", "-3.6%", "of GDP  (deficit)", RED_WARN),
    ("Current Account\n2021", "+11.9%", "of GDP  (surplus)", ACCENT2),
    ("IMF ECF\nProgramme", "$1.3 B", "approved 2022", ACCENT),
]
card_left = Inches(7.3)
card_top  = Inches(1.8)
card_w    = Inches(1.85)
card_h    = Inches(1.6)
card_gap  = Inches(0.15)
for i, (m, v, s, vc) in enumerate(cards):
    stat_card(s3,
              card_left + i * (card_w + card_gap),
              card_top, card_w, card_h, m, v, s, vc)

# Current Account column chart — matches PBI clusteredColumnChart on "Copper Triggered the Recovery"
add_textbox(s3, Inches(7.3), Inches(3.65), Inches(5.85), Inches(0.3),
            "CURRENT ACCOUNT BALANCE  (% of GDP)  ·  2015–2024", 8, bold=True, color=MID_TEXT)
add_white_bg(s3, Inches(7.3), Inches(3.95), Inches(5.85), Inches(2.75))
add_column_chart(s3, Inches(7.3), Inches(3.95), Inches(5.85), Inches(2.75),
                 YEARS, CURR_ACCT, series_color=ACCENT2,
                 point_colors={0: RED_WARN, 1: RED_WARN, 2: RED_WARN,
                                3: RED_WARN, 4: RED_WARN,
                                6: ACCENT2, 7: ACCENT2, 8: ACCENT2, 9: ACCENT2})

add_textbox(s3, Inches(12.5), Inches(6.9), Inches(0.7), Inches(0.4),
            "03", 9, color=MID_TEXT, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — The Structural Fix: G20 Common Framework
# ══════════════════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(blank_layout)
set_bg(s4, DARK_BG)
add_rect(s4, Inches(0), Inches(0), Inches(0.12), SLIDE_H, ACCENT2)
add_rect(s4, Inches(0.12), Inches(0), SLIDE_W - Inches(0.12), Inches(0.06), DIVIDER)

add_textbox(s4, Inches(0.5), Inches(0.25), Inches(9), Inches(0.4),
            "3-MINUTE STORY  ·  PART 3 OF 4  ·  THE STRUCTURAL FIX", 9, color=ACCENT2, bold=True)

add_multiline(s4, Inches(0.5), Inches(0.7), Inches(12), Inches(1.1),
              [("Debt Restructuring Locked In the Gains", True, LIGHT_TEXT)],
              font_size=28)

add_divider_line(s4, Inches(0.5), Inches(1.75), Inches(5.5), color=ACCENT2)

narrative4 = [
    ("In November 2023, Zambia finalised its G20 Common Framework restructuring deal — "
     "the first successful completion under this framework. Bilateral creditor agreements "
     "were implemented through 2024.", False, LIGHT_TEXT),
    ("", False, LIGHT_TEXT),
    ("The structural impact was immediate and measurable:", False, MID_TEXT),
    ("  •  PPG debt service fell from 12.5 % → 3 % of exports", False, ACCENT2),
    ("  •  GDP growth held at 5.2–5.4 % even as copper prices cooled", False, ACCENT2),
    ("  •  FDI inflows hit a decade-high of 9.32 % of GDP in 2024", False, ACCENT2),
    ("", False, LIGHT_TEXT),
    ("This is the key distinction: copper created fiscal space; restructuring "
     "institutionalised it.", True, LIGHT_TEXT),
]
add_multiline(s4, Inches(0.5), Inches(2.0), Inches(6.2), Inches(3.8),
              narrative4, font_size=12)

# before / after comparison
add_textbox(s4, Inches(7.3), Inches(1.8), Inches(5.8), Inches(0.35),
            "BEFORE vs AFTER RESTRUCTURING", 9, bold=True, color=MID_TEXT)

rows = [
    ("PPG Debt Service / Exports", "12.5 %", "3.0 %"),
    ("GDP Growth", "-2.8 %  (2020)", "5.3 %  (2024)"),
    ("FDI Inflows (% GDP)", "1.8 %  (2020)", "9.32 %  (2024)"),
    ("Investor Confidence", "Eurobond default", "Decade-high FDI"),
]
row_top = Inches(2.25)
for i, (metric, before, after) in enumerate(rows):
    y = row_top + i * Inches(0.92)
    add_rect(s4, Inches(7.3), y, Inches(5.85), Inches(0.82), DIVIDER)
    add_textbox(s4, Inches(7.4), y + Inches(0.05), Inches(2.3), Inches(0.35),
                metric, 9, color=MID_TEXT)
    add_textbox(s4, Inches(7.4), y + Inches(0.4), Inches(2.0), Inches(0.35),
                before, 12, bold=True, color=RED_WARN)
    add_textbox(s4, Inches(10.2), y + Inches(0.4), Inches(2.7), Inches(0.35),
                "→  " + after, 12, bold=True, color=ACCENT2)

add_textbox(s4, Inches(12.5), Inches(6.9), Inches(0.7), Inches(0.4),
            "04", 9, color=MID_TEXT, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — The Evidence: Key Metrics at a Glance
# ══════════════════════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(blank_layout)
set_bg(s5, DARK_BG)
add_rect(s5, Inches(0), Inches(0), Inches(0.12), SLIDE_H, ACCENT)
add_rect(s5, Inches(0.12), Inches(0), SLIDE_W - Inches(0.12), Inches(0.06), DIVIDER)

add_textbox(s5, Inches(0.5), Inches(0.25), Inches(9), Inches(0.4),
            "3-MINUTE STORY  ·  PART 4 OF 4  ·  THE EVIDENCE", 9, color=ACCENT, bold=True)

add_multiline(s5, Inches(0.5), Inches(0.7), Inches(12), Inches(1.1),
              [("The Numbers Confirm the Verdict", True, LIGHT_TEXT)],
              font_size=28)

add_divider_line(s5, Inches(0.5), Inches(1.75), Inches(12.5))

# 6 KPI cards in 2 rows
kpis = [
    ("External Debt Peak", "$29.1 B", "2019  (tripled in 4 yrs)", RED_WARN),
    ("Current Account Swing", "+15.5 pp", "from -3.6% to +11.9% GDP", ACCENT2),
    ("IMF ECF Approved", "$1.3 B", "2022 facility", ACCENT),
    ("PPG Debt Service\n(post-deal)", "3.0 %", "down from 12.5% of exports", ACCENT2),
    ("GDP Growth (2024)", "5.3 %", "sustained as Cu prices cooled", ACCENT2),
    ("FDI Inflows (2024)", "9.32 %", "decade high, % of GDP", ACCENT),
]
kpi_w = Inches(2.0)
kpi_h = Inches(1.55)
kpi_gap_x = Inches(0.22)
kpi_gap_y = Inches(0.3)
kpi_start_x = Inches(0.5)
kpi_start_y = Inches(2.05)

for i, (m, v, s, vc) in enumerate(kpis):
    col = i % 3
    row = i // 3
    x = kpi_start_x + col * (kpi_w + kpi_gap_x)
    y = kpi_start_y + row * (kpi_h + kpi_gap_y)
    # wider cards — 3 per row but across full width
    w = Inches(4.1)
    x = Inches(0.5) + col * (Inches(4.1) + Inches(0.12))
    stat_card(s5, x, y, w, kpi_h, m, v, s, vc)

# Data sources note
add_textbox(s5, Inches(0.5), Inches(6.55), Inches(12), Inches(0.55),
            "Sources: World Bank WDI · IMF World Economic Outlook · Bank of Zambia · "
            "Ministry of Finance — Report on the Economy 2025",
            8, color=MID_TEXT)

add_textbox(s5, Inches(12.5), Inches(6.9), Inches(0.7), Inches(0.4),
            "05", 9, color=MID_TEXT, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Charts: GDP Growth + FDI Inflows (2015–2024)
# ══════════════════════════════════════════════════════════════════════════════
s6c = prs.slides.add_slide(blank_layout)
set_bg(s6c, DARK_BG)
add_rect(s6c, Inches(0), Inches(0), Inches(0.12), SLIDE_H, ACCENT)
add_rect(s6c, Inches(0.12), Inches(0), SLIDE_W - Inches(0.12), Inches(0.06), DIVIDER)

add_textbox(s6c, Inches(0.5), Inches(0.25), Inches(9), Inches(0.4),
            "DATA EVIDENCE  ·  CHARTS", 9, color=ACCENT, bold=True)
add_multiline(s6c, Inches(0.5), Inches(0.65), Inches(12.5), Inches(0.8),
              [("GDP Growth & FDI Inflows: The Recovery Visualised", True, LIGHT_TEXT)],
              font_size=22)
add_divider_line(s6c, Inches(0.5), Inches(1.48), Inches(12.5))

# Explanatory intro sentence
add_textbox(s6c, Inches(0.5), Inches(1.58), Inches(12.3), Inches(0.42),
            "GDP contracted sharply in 2020 but rebounded within a year — and sustained 5 %+ growth "
            "even as copper prices cooled post-2021. FDI hit its lowest point at the 2020 default, "
            "then climbed to a decade-high 9.32 % in 2024 — the year bilateral restructuring "
            "agreements were fully implemented, confirming restored investor confidence.",
            9, color=MID_TEXT)

# GDP Growth chart (left half)
add_textbox(s6c, Inches(0.5), Inches(2.08), Inches(6.2), Inches(0.28),
            "GDP GROWTH  (% annual)  ·  2015–2024  ·  Red = contraction  |  Green = recovery",
            8, bold=True, color=MID_TEXT)
add_white_bg(s6c, Inches(0.5), Inches(2.38), Inches(6.2), Inches(3.8))
add_column_chart(s6c, Inches(0.5), Inches(2.38), Inches(6.2), Inches(3.8),
                 YEARS, GDP_GROWTH, series_color=ACCENT,
                 point_colors={4: RGBColor(0xFF, 0xA5, 0x00),
                                5: RED_WARN,
                                6: ACCENT2, 7: ACCENT2, 8: ACCENT2, 9: ACCENT2})
add_textbox(s6c, Inches(0.5), Inches(6.2), Inches(6.2), Inches(0.32),
            "GDP fell -2.8 % in 2020 (COVID + default) before recovering to 4.6 % in 2021 "
            "and holding at 5.3 % in 2024 despite softer copper prices.",
            8, color=MID_TEXT, italic=True)

# FDI chart (right half)
add_textbox(s6c, Inches(7.0), Inches(2.08), Inches(6.0), Inches(0.28),
            "FDI NET INFLOWS  (% of GDP)  ·  2015–2024  ·  Red = decade-low  |  Green = decade-high",
            8, bold=True, color=MID_TEXT)
add_white_bg(s6c, Inches(7.0), Inches(2.38), Inches(6.0), Inches(3.8))
add_column_chart(s6c, Inches(7.0), Inches(2.38), Inches(6.0), Inches(3.8),
                 YEARS, FDI, series_color=ACCENT,
                 point_colors={5: RED_WARN, 9: ACCENT2})
add_textbox(s6c, Inches(7.0), Inches(6.2), Inches(6.0), Inches(0.32),
            "FDI bottomed at 1.9 % in 2020 then surged to 9.32 % in 2024 — "
            "the year G20 bilateral restructuring agreements were implemented.",
            8, color=MID_TEXT, italic=True)

add_textbox(s6c, Inches(0.5), Inches(6.82), Inches(12), Inches(0.4),
            "Sources: World Bank WDI · IMF WEO · Bank of Zambia · Ministry of Finance — Report on the Economy 2025",
            8, color=MID_TEXT)
add_textbox(s6c, Inches(12.5), Inches(6.9), Inches(0.7), Inches(0.4),
            "06", 9, color=MID_TEXT, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — The Verdict + Call to Action
# ══════════════════════════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(blank_layout)
set_bg(s6, DARK_BG)
add_rect(s6, Inches(0), Inches(0), Inches(0.12), SLIDE_H, ACCENT2)
add_rect(s6, Inches(0.12), Inches(0), SLIDE_W - Inches(0.12), Inches(0.06), DIVIDER)

add_textbox(s6, Inches(0.5), Inches(0.25), Inches(9), Inches(0.4),
            "CONCLUSION  ·  THE VERDICT", 9, color=ACCENT2, bold=True)

add_multiline(s6, Inches(0.5), Inches(0.7), Inches(12.5), Inches(1.2),
              [("Both Were Necessary — Neither Was Sufficient Alone", True, LIGHT_TEXT)],
              font_size=26)

add_divider_line(s6, Inches(0.5), Inches(1.85), Inches(12.3), color=ACCENT2)

# Two-column verdict
left_head  = "COPPER PROVIDED THE WINDFALL"
left_body  = (
    "The global copper price supercycle of 2021 handed Zambia a "
    "fiscal windfall: the current account swung by +15.5 percentage "
    "points and exports surged. This created the fiscal space that "
    "made restructuring politically viable and economically credible.\n\n"
    "Without the copper boom, there would have been nothing to restructure "
    "toward — and no creditors willing to take a haircut on a country "
    "with no revenue trajectory."
)
right_head = "RESTRUCTURING MADE IT DURABLE"
right_body = (
    "The G20 Common Framework deal converted a temporary commodity "
    "windfall into a permanent reduction in debt service obligations. "
    "PPG debt service fell from 12.5 % to 3 % of exports — freeing "
    "fiscal resources for productive investment.\n\n"
    "GDP growth of 5.2–5.4 % persisted even as copper prices moderated, "
    "and FDI inflows hit a decade-high of 9.32 % of GDP — confirming "
    "that investors saw institutional, not just commodity, credibility."
)

for side, head, body, col in [
    (0, left_head,  left_body,  ACCENT),
    (1, right_head, right_body, ACCENT2),
]:
    x = Inches(0.5) + side * Inches(6.5)
    add_textbox(s6, x, Inches(2.1), Inches(6.0), Inches(0.45),
                head, 11, bold=True, color=col)
    add_textbox(s6, x, Inches(2.6), Inches(6.15), Inches(3.0),
                body, 11, color=LIGHT_TEXT)

# vertical separator
add_rect(s6, Inches(6.55), Inches(2.0), Inches(0.04), Inches(3.6), DIVIDER)

# Final Big Idea reprise
add_rect(s6, Inches(0.5), Inches(5.8), Inches(12.3), Inches(1.05), DIVIDER)
add_textbox(s6, Inches(0.65), Inches(5.87), Inches(1.2), Inches(0.32),
            "BIG IDEA", 8, bold=True, color=ACCENT)
big_idea_short = (
    "Zambia's recovery proves that commodity luck and institutional reform are complements, "
    "not substitutes — copper created the window; the G20 deal made it a door."
)
add_textbox(s6, Inches(0.65), Inches(6.17), Inches(12.0), Inches(0.55),
            big_idea_short, 13, bold=True, color=LIGHT_TEXT, italic=True)

add_textbox(s6, Inches(12.5), Inches(6.9), Inches(0.7), Inches(0.4),
            "07", 9, color=MID_TEXT, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Appendix: 3-Minute Story & Big Idea (text version)
# ══════════════════════════════════════════════════════════════════════════════
s7 = prs.slides.add_slide(blank_layout)
set_bg(s7, DARK_BG)
add_rect(s7, Inches(0), Inches(0), Inches(0.12), SLIDE_H, MID_TEXT)
add_rect(s7, Inches(0.12), Inches(0), SLIDE_W - Inches(0.12), Inches(0.06), DIVIDER)

add_textbox(s7, Inches(0.5), Inches(0.25), Inches(9), Inches(0.4),
            "APPENDIX  ·  3-MINUTE STORY & BIG IDEA", 9, color=MID_TEXT, bold=True)

add_multiline(s7, Inches(0.5), Inches(0.65), Inches(12), Inches(0.7),
              [("Storytelling with Data — Written Narrative", True, LIGHT_TEXT)],
              font_size=20)

add_divider_line(s7, Inches(0.5), Inches(1.35), Inches(12.3), color=MID_TEXT)

three_min = (
    "Between 2015 and 2019, Zambia borrowed heavily against future copper revenues, "
    "tripling its external debt from $12.3 B to $29.1 B. When COVID-19 struck in 2020, "
    "revenues collapsed and Zambia became the first African sovereign pandemic-era default. "
    "The question we needed to answer: was the recovery that followed real — and sustainable?\n\n"
    "In 2021, a global copper price supercycle delivered a fiscal windfall. The current "
    "account swung 15.5 percentage points — from a -3.6 % deficit to a +11.9 % surplus — "
    "before the restructuring deal had even been negotiated. Copper created the breathing room.\n\n"
    "But breathing room is not the same as structural reform. In November 2023, Zambia "
    "finalised the G20 Common Framework deal, becoming the first country to complete the "
    "process. The results were clear: PPG debt service fell from 12.5 % to 3 % of exports; "
    "GDP growth held at 5.2–5.4 % even as copper prices cooled; and FDI inflows reached "
    "a decade-high 9.32 % of GDP in 2024.\n\n"
    "The verdict: both were necessary. Commodity luck created the window. Institutional "
    "reform made it a door. Going forward, Zambia's challenge is to diversify beyond copper "
    "— into green energy, agriculture, and manufacturing — before the next commodity "
    "downcycle tests the depth of these institutional gains."
)
add_textbox(s7, Inches(0.5), Inches(1.5), Inches(6.2), Inches(5.2),
            three_min, 10.5, color=LIGHT_TEXT)

# Big Idea box
add_rect(s7, Inches(7.0), Inches(1.5), Inches(5.9), Inches(2.4), DIVIDER)
add_textbox(s7, Inches(7.15), Inches(1.6), Inches(3), Inches(0.35),
            "BIG IDEA", 9, bold=True, color=ACCENT)
big_idea_full = (
    "Zambia's G20 Common Framework debt restructuring converted a copper-driven "
    "fiscal windfall into a durable institutional framework — slashing PPG debt "
    "service from 12.5 % to 3 % of exports, sustaining 5 %+ GDP growth, and "
    "attracting record FDI of 9.32 % of GDP in 2024 — proving that commodity "
    "luck alone cannot anchor a recovery without credible institutional reform."
)
add_textbox(s7, Inches(7.15), Inches(2.0), Inches(5.6), Inches(1.75),
            big_idea_full, 11, color=LIGHT_TEXT, italic=True)

# Big Idea checklist
checklist = [
    ("✓  Unique point of view", ACCENT2),
    ("✓  Conveys what's at stake", ACCENT2),
    ("✓  Complete sentence", ACCENT2),
]
for i, (txt, col) in enumerate(checklist):
    add_textbox(s7, Inches(7.15), Inches(4.05) + i * Inches(0.38),
                Inches(5.6), Inches(0.35), txt, 10, color=col, bold=True)

add_textbox(s7, Inches(12.5), Inches(6.9), Inches(0.7), Inches(0.4),
            "08", 9, color=MID_TEXT, align=PP_ALIGN.RIGHT)


# ── Save ──────────────────────────────────────────────────────────────────────
out_path = r"c:\Users\bmweemba\Documents\Boldwin-Mweemba-Zambia Economic recovery\Zambia_Economic_Recovery_Presentation.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
